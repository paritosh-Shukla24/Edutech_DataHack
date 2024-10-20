from fastapi import FastAPI, File, UploadFile, HTTPException
import uvicorn
from PIL import Image,PngImagePlugin
import google.generativeai as genai
import json
import ast
import os
from dotenv import load_dotenv
from pydub import AudioSegment
from deepGram import Transcriber
from langchain_community.tools import YouTubeSearchTool
from io import BytesIO
import pyttsx3
from fastapi import FastAPI, Depends, HTTPException, status
from pydantic import BaseModel, EmailStr

from dotenv import load_dotenv
import os
from pydantic import BaseModel
from Ai import QuizRLAgent
app = FastAPI()
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
quiz=QuizRLAgent(initial_score=50)
genai.configure(api_key=api_key)
model_vision = genai.GenerativeModel('gemini-1.5-flash-latest')
model_text=genai.GenerativeModel('gemini-1.5-pro')
engine = pyttsx3.init()
voices = engine.getProperty('voices')
engine.setProperty('voice', voices[1].id)  
engine.setProperty('rate', 150)  
engine.setProperty('volume', 0.9)  
tool=YouTubeSearchTool()
from fastapi import FastAPI, HTTPException, Depends
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from pydantic import BaseModel
from bson import ObjectId
from pymongo import MongoClient
from passlib.context import CryptContext
from jose import JWTError, jwt
from datetime import datetime, timedelta
import pytesseract
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
difficulty="easy"
SECRET_KEY = "your_secret_key"
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30
# MongoDB setup
client = MongoClient('mongodb+srv://Hackathon:59ILYisG4iNaGNRU@cluster0.w6oegog.mongodb.net/?retryWrites=true&w=majority&appName=Cluster0')
db = client.get_database('Hackathon')  # Replace with your DB name
users_collection = db.get_collection('users')
scores_collection = db.get_collection('gmail_scores')
# FastAPI app
app = FastAPI()
from fastapi.middleware.cors import CORSMiddleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Replace with your frontend URL
    allow_credentials=True,
    allow_methods=["*"],  # Allow all HTTP methods
    allow_headers=["*"],  # Allow all headers
)
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
# Pydantic models
class Token(BaseModel):
    access_token: str
    token_type: str

class TokenData(BaseModel):
    email: str

class User(BaseModel):
    email: str
    password: str

# JWT utility functions
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

def create_access_token(data: dict, expires_delta: timedelta = None):
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.utcnow() + expires_delta
    else:
        expire = datetime.utcnow() + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    return pwd_context.hash(password)

# API endpoints
@app.post("/login")
async def login(form_data: OAuth2PasswordRequestForm = Depends()):
    user = users_collection.find_one({"email": form_data.username})
    
    if not user or not verify_password(form_data.password, user['password']):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Invalid credentials",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    access_token = create_access_token(data={"email": user['email']})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/register")
async def register(user: User):
    if users_collection.find_one({"email": user.email}):
        raise HTTPException(status_code=400, detail="Email already registered")
    
    hashed_password = get_password_hash(user.password)
    users_collection.insert_one({"email": user.email, "password": hashed_password})
    return {"message": "User registered successfully"}

@app.get("/me")
async def read_users_me(token: str = Depends(oauth2_scheme)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("email")
        if email is None:
            raise credentials_exception
        return {"email": email}
    except JWTError:
        raise credentials_exception




def Answering(question,answer,options,format):
    if format=='text':
        input_prompt=f"""Yoy are Question Solver which solve the Question and find the appropriate solution.
        Question : {question}
        Options : {options}
        Give detailed solution in 50 words max
        """
        response=model_text.generate_content(input_prompt)
        return response.text
    if format=="audio":
        input_prompt=f"""Yoy are Question Solver which solve the Question and find the appropriate solution and explain solution.
        Question : {question}
        Options : {options}
        Give detailed solution in 100 words max
        """
        response=model_text.generate_content(input_prompt)
        engine.save_to_file(response.text, "output.wav")
        engine.runAndWait()
    if format=="video":
        link=tool.run(answer)
        return link
    
def extract_text_from_pdf(file_content):
    reader = PdfReader(BytesIO(file_content))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_content):
    doc = Document(BytesIO(file_content))
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_pptx(file_content):
    presentation = Presentation(BytesIO(file_content))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(image_content):
    # Use pytesseract to extract text from an image
    image = Image.open(BytesIO(image_content))
    text = pytesseract.image_to_string(image)
    return text


global_image = None  
global_text_audio=None

@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    content_type = file.content_type
    global global_image  
    print(content_type)
    global difficulty
    
    if content_type in ["application/pdf", "image/jpeg", "image/png", "audio/mpeg", "video/mp4", 
                   "application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        file_content = await file.read()  
        if content_type == "application/pdf":
            
            text = extract_text_from_pdf(file_content)
            txt_file_path = f"./extracted_text.txt"
            with open(txt_file_path, "w",encoding='utf-8') as txt_file:
                txt_file.write(text)
            previous = "What is the main topic discussed in this content?"
            input_prompt = f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the content provided. 
             Content: {text}
        Difficulty of Question is: {difficulty}
        Please look at the previous Question and make sure you don't continue the Question.
        Previous Question: {previous}
        """
            prompt = """
        Generate 1 question and return in JSON format exactly 1 question with 4 options and also the correct answer and a small solution, with JSON attributes named question, options, answer, solution.
        """
            response = model_vision.generate_content([input_prompt, prompt])
            clean_response = response.text.strip().replace("```json", "").replace("```", "")
            #print(clean_response)
            #print(type(clean_response))
            quiz_data = json.loads(clean_response)
            format='video'
            link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
            link_list = ast.literal_eval(link)
            first_link = link_list[0]
            return {
            "question": quiz_data['question'],
            "options": quiz_data['options'],
            "answer": quiz_data['answer'],
            "solution": quiz_data['solution'],
            "link":first_link
        }
    

        elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = extract_text_from_docx(file_content)
            txt_file_path = f"./extracted_text.txt"
            with open(txt_file_path, "w",encoding='utf-8') as txt_file:
                txt_file.write(text)
            
            previous = "What is the main topic discussed in this content?"
            input_prompt = f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the content provided. 
             Content: {text}
        Difficulty of Question is: {difficulty}
        Please look at the previous Question and make sure you don't continue the Question.
        Previous Question: {previous}
        """
            prompt = """
        Generate 1 question and return in JSON format exactly 1 question with 4 options and also the correct answer and a small solution, with JSON attributes named question, options, answer, solution.
        """
            response = model_vision.generate_content([input_prompt, prompt])
            clean_response = response.text.strip().replace("json", "").replace("", "")
            print(clean_response)
            print(type(clean_response))
            quiz_data = json.loads(clean_response)
            format='video'
            link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
            link_list = ast.literal_eval(link)
            first_link = link_list[0]

            return {
            "question": quiz_data['question'],
            "options": quiz_data['options'],
            "answer": quiz_data['answer'],
            "solution": quiz_data['solution'],
            "link":first_link
        }
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = extract_text_from_pptx(file_content)
            txt_file_path = f"./extracted_text.txt"
            with open(txt_file_path, "w",encoding='utf-8') as txt_file:
                txt_file.write(text)
            previous = "What is the main topic discussed in this content?"
            input_prompt = f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the content provided. 
             Content: {text}
        Difficulty of Question is: {difficulty}
        Please look at the previous Question and make sure you don't continue the Question.
        Previous Question: {previous}
        """
            prompt = """
        Generate 1 question and return in JSON format exactly 1 question with 4 options and also the correct answer and a small solution, with JSON attributes named question, options, answer, solution.
        """
            response = model_vision.generate_content([input_prompt, prompt])
            clean_response = response.text.strip().replace("json", "").replace("", "")
            print(clean_response)
            print(type(clean_response))
            quiz_data = json.loads(clean_response)
            format='video'
            link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
            link_list = ast.literal_eval(link)
            first_link = link_list[0]

            return {
            "question": quiz_data['question'],
            "options": quiz_data['options'],
            "answer": quiz_data['answer'],
            "solution": quiz_data['solution'],
            "link":first_link
        }

        elif content_type in ["image/jpeg", "image/png"]:
            image = Image.open(BytesIO(file_content))
            os.makedirs("./images", exist_ok=True)
            image_path = f"./images/saved_image.png"  # Save to a proper directory
            image.save(image_path)  #
            previous="What is the direction of the force acting on the second object, m2, in relation to the direction of the force acting on the first object, m1"
            input_prompt=f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the content of the image provided. 
            Difficulty of Question is : {difficulty}
            Please look at previous Question make sure you don't continue the Question.
            previous Question : {previous}
            """
            prompt="""
            generate 1 Questions and return in json format exactly 1 questions with 4 options and and also correct answer and small solution,with json attributes named question,options,answer,solution.
            """
            response = model_vision.generate_content([input_prompt,image, prompt])
            clean_response = response.text.strip().replace("```json", "").replace("```", "")
            quiz_data = json.loads(clean_response)
            #print(quiz_data)
            format='video'
            link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
            link_list = ast.literal_eval(link)
            first_link = link_list[0]
            return {
                "question":quiz_data['question'],
                "options":quiz_data['options'],
                "answer":quiz_data['answer'],
                "solution": quiz_data['solution'],
                "link":first_link
            }
        elif content_type == "audio/mpeg":
                audio = AudioSegment.from_file(BytesIO(file_content), format="mp3")
                wav_filename = f"temp.wav"
                audio.export(wav_filename, format="wav")
                text=Transcriber(wav_filename)
                previous='Which of the following is NOT a common type of neural network used in deep learning?'
                input_prompt=f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the User domain which he is aking. Based on Topic  Generate.
            Topic : {text} 
            Difficulty of Question is : {difficulty}
           
            Previous Question : {previous}
             Please look at previous Question make sure you don't repeat and try new type of questions.Completely different Question.

            generate 1 Questions and return in json format exactly 1 questions with 4 options and and also correct answer and small solution,with json attributes named question,options,answer,solution.
            """
                response = model_vision.generate_content(input_prompt)
                clean_response = response.text.strip().replace("```json", "").replace("```", "")
                quiz_data = json.loads(clean_response)
                format='video'
                link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
                link_list = ast.literal_eval(link)
                first_link = link_list[0]
                return {
                "question":quiz_data['question'],
                "options":quiz_data['options'],
                "answer":quiz_data['answer'],
                "solution":quiz_data['solution'],
                "link": first_link
                }
       

        elif content_type == "video/mp4":
            return {"filename": file.filename, "type": "Video file received"}
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    return {"filename": file.filename, "type": "Unknown file type"}

class NextQuestionRequest(BaseModel):
    type: str  

@app.post("/next_questions/")
async def next_question(request_data: NextQuestionRequest):
    content_type = request_data.type.lower()
    global difficulty
    if content_type == "pdf":
        file_path= f"./extracted_text.txt"
        with open(file_path, "r", encoding="utf-8") as file:
            file_content = file.read()
        #print(file_content)    

        
        previous = "What is the main topic discussed in this content?"
        input_prompt = f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the content provided. 
             Content: {file_content}
        Difficulty of Question is: {difficulty}
        Please look at the previous Question and make sure you don't continue the Question.
        Previous Question: {previous}
        """
        prompt = """
        Generate 1 question and return in JSON format exactly 1 question with 4 options and also the correct answer and a small solution, with JSON attributes named question, options, answer, solution.
        """
        response = model_vision.generate_content([input_prompt, prompt])
        clean_response = response.text.strip().replace("```json", "").replace("```", "")
        #print(clean_response)
        #print(type(clean_response))
        quiz_data = json.loads(clean_response)
        format='video'
        link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
        link_list = ast.literal_eval(link)
        first_link = link_list[0]
        return {
            "question": quiz_data['question'],
            "options": quiz_data['options'],
            "answer": quiz_data['answer'],
            "solution": quiz_data['solution'],
            "link":first_link
        }    
            
    elif content_type in ["jpeg", "png"]:
            image_path = "./images/saved_image.png"
            image = Image.open(image_path)
            previous="What is the direction of the force acting on the second object, m2, in relation to the direction of the force acting on the first object, m1"
            input_prompt=f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the content of the image provided. 
            Difficulty of Question is : {difficulty}
            Please look at previous Question make sure you don't continue the Question.
            previous Question : {previous}
            """
            prompt="""
            generate 1 Questions and return in json format exactly 1 questions with 4 options and and also correct answer and small solution,with json attributes named question,options,answer,solution.
            """
            response = model_vision.generate_content([input_prompt,image, prompt])
            clean_response = response.text.strip().replace("```json", "").replace("```", "")
            quiz_data = json.loads(clean_response)
            format='video'
            link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
            link_list = ast.literal_eval(link)
            first_link = link_list[0]
            return {
                "question":quiz_data['question'],
                "options":quiz_data['options'],
                "answer":quiz_data['answer'],
                "solution": quiz_data['solution'],
                "link":first_link
            }
    elif content_type == "mp3":
            file_path=r'temp.wav'
            audio = AudioSegment.from_wav(file_path)
            global_text_audio=Transcriber(audio)
            
            previous='Which of the following is NOT a common type of neural network used in deep learning?'
            input_prompt=f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the User domain which he is aking. Based on Topic  Generate.
            Topic : {global_text_audio} 
            Difficulty of Question is : {difficulty}
           
            Previous Question : {previous}
             Please look at previous Question make sure you don't repeat and try new type of questions.Completely different Question.

            generate 1 Questions and return in json format exactly 1 questions with 4 options and and also correct answer and small solution,with json attributes named question,options,answer,solution.
            """
            response = model_vision.generate_content(input_prompt)
            clean_response = response.text.strip().replace("```json", "").replace("```", "")
            quiz_data = json.loads(clean_response)
            format='video'
            link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
            print(link)
            link_list = ast.literal_eval(link)
            first_link = link_list[0]
            return {
                "question":quiz_data['question'],
                "options":quiz_data['options'],
                "answer":quiz_data['answer'],
                "solution": quiz_data['solution'],
                "link ": first_link
                }
    elif content_type == "docx":
        txt_file_path = f"./extracted_text.txt"
        with open(txt_file_path, "r", encoding="utf-8") as file:
            file_content = file.read()
        
        previous = "What is the main topic discussed in this content?"
        input_prompt = f"""You are a Quiz Generator AI. 
        Your task is to generate a set of quiz questions based on the content provided. 
             Content: {file_content}
        Difficulty of Question is: {difficulty}
        Please look at the previous Question and make sure you don't continue the Question.
        Previous Question: {previous}
        """
        prompt = """
        Generate 1 question and return in JSON format exactly 1 question with 4 options and also the correct answer and a small solution, with JSON attributes named question, options, answer, solution.
        """
        response = model_vision.generate_content([input_prompt, prompt])
        clean_response = response.text.strip().replace("json", "").replace("", "")
        print(clean_response)
        print(type(clean_response))
        quiz_data = json.loads(clean_response)
        format='video'
        link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
        link_list = ast.literal_eval(link)
        first_link = link_list[0]
        return {
            "question": quiz_data['question'],
            "options": quiz_data['options'],
            "answer": quiz_data['answer'],
            "solution": quiz_data['solution'],
            "link":first_link
        }    
    elif content_type=="pptx":
            txt_file_path = f"./extracted_text.txt"
            with open(txt_file_path, "r", encoding="utf-8") as file:
                file_content = file.read()
          
            previous = "What is the main topic discussed in this content?"
            input_prompt = f"""You are a Quiz Generator AI. 
            Your task is to generate a set of quiz questions based on the content provided. 
             Content: {file_content}
        Difficulty of Question is: {difficulty}
        Please look at the previous Question and make sure you don't continue the Question.
        Previous Question: {previous}
        """
            prompt = """
        Generate 1 question and return in JSON format exactly 1 question with 4 options and also the correct answer and a small solution, with JSON attributes named question, options, answer, solution.
        """
            response = model_vision.generate_content([input_prompt, prompt])
            clean_response = response.text.strip().replace("json", "").replace("", "")
            print(clean_response)
            print(type(clean_response))
            quiz_data = json.loads(clean_response)
            format='video'
            link=Answering(quiz_data['question'],quiz_data['answer'],quiz_data['options'],format)
            link_list = ast.literal_eval(link)
            first_link = link_list[0]
            return {
            "question": quiz_data['question'],
            "options": quiz_data['options'],
            "answer": quiz_data['answer'],
            "solution": quiz_data['solution'],
            "link":first_link
        }
    elif content_type == "mp4":
            return {"filename": file.filename, "type": "Video file received"}
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")


class GmailScoreModel(BaseModel):
    gmail: str
    result: bool 
    difficulty: str
    response_time: float

@app.post("/update_gmail_score/")
def update_gmail_score(gmail_score: GmailScoreModel):
    current_time = datetime.utcnow()
    existing_entry = scores_collection.find_one({"gmail": gmail_score.gmail})

    if existing_entry:
        # Extract the initial score from the existing entry's history or score field
        initial_score = existing_entry.get('history', [{}])[-1].get('score', 70)  # Default score is 70 if not found
        quiz = QuizRLAgent(initial_score=initial_score)

        print(gmail_score.difficulty)
        print(gmail_score.response_time)
        print(gmail_score.result)
        
        
        gmail_score.response_time = int(gmail_score.response_time)  # This may be redundant since it's a float

        # Adjust score using the QuizRLAgent
        updated_score = quiz.adjust_score(gmail_score.difficulty, gmail_score.response_time, gmail_score.result)
        if updated_score < 80 and updated_score >70:
            global difficulty
            difficulty="medium"
        if updated_score <100 and updated_score>80:
            difficulty="hard"    
        update_entry = {
            "score": updated_score,
            "timestamp": current_time
        }

        # Update the database with the new score and timestamp
        scores_collection.update_one(
            {"gmail": gmail_score.gmail},
            {
                "$push": {"history": update_entry},
                "$set": {"last_updated": current_time}
            }
        )
        return {"message": "Score updated", "gmail": gmail_score.gmail, "last_update": update_entry}
    
    else:
        # Handle the case when there is no existing entry for the given gmail
        new_entry = {
            "gmail": gmail_score.gmail,
            "last_updated": current_time,
            "history": [{
                "score": 70,  # Start with the default score of 70
                "timestamp": current_time
            }]
        }
        result = scores_collection.insert_one(new_entry)
        return {"message": "Gmail and score added", "id": str(result.inserted_id), "last_update": new_entry["history"][-1]}
from typing import Optional
@app.get("/get_gmail_scores/")
def get_gmail_scores(gmail: str, date: Optional[str] = None):
    existing_entry = scores_collection.find_one({"gmail": gmail})
    if not existing_entry:
        return {"error": "Gmail not found."}

    history = existing_entry.get("history", [])
    
    if date:
        # Parse the date string
        try:
            date_obj = datetime.strptime(date, "%Y-%m-%d")
        except ValueError:
            return {"error": "Invalid date format. Use YYYY-MM-DD."}
        
        # Filter history entries where timestamp is on the given date
        filtered_history = []
        for entry in history:
            entry_timestamp = entry.get("timestamp")
            if isinstance(entry_timestamp, str):
                # Convert timestamp string to datetime object
                entry_timestamp = datetime.fromisoformat(entry_timestamp)
            if entry_timestamp.date() == date_obj.date():
                filtered_history.append(entry)   

        return {"history": filtered_history}
    else:
        return { "history": history}
    

import networkx as nx
from langchain.agents import AgentExecutor
from langchain_community.tools.tavily_search import TavilySearchResults
from langchain_groq import ChatGroq
from langchain import hub
from langchain.agents import create_tool_calling_agent
import matplotlib.pyplot as plt
from fastapi.responses import FileResponse
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
TAVILY_API_KEY = os.getenv("TAVILY_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
LANGCHAIN_API_KEY = os.getenv("LANGCHAIN_API_KEY")
LANGCHAIN_PROJECT = os.getenv("LANGCHAIN_PROJECT")

os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
os.environ["TAVILY_API_KEY"] = TAVILY_API_KEY
os.environ["GROQ_API_KEY"] = GROQ_API_KEY
os.environ["LANGCHAIN_API_KEY"] = LANGCHAIN_API_KEY
os.environ["LANGCHAIN_TRACING_V2"] = "true"
os.environ["LANGCHAIN_ENDPOINT"] = "https://api.smith.langchain.com"
os.environ["LANGCHAIN_PROJECT"] = LANGCHAIN_PROJECT

search = TavilySearchResults()
tools = [search]
llm3 = ChatGroq(model_name="llama3-70b-8192")


prompt = hub.pull("hwchase17/openai-functions-agent")
agent = create_tool_calling_agent(llm3, tools, prompt)
agent_executor = AgentExecutor(agent=agent, tools=tools, verbose=True)
class DomainRequest(BaseModel):
    domain_name: str
@app.post("/generate-roadmap/")
async def generate_roadmap(request: DomainRequest):
    domain_name = request.domain_name
    # Generate roadmap and course recommendations
    combined_prompt = f"""
    Generate a comprehensive 8 to 10-step learning roadmap for {domain_name}. Each step should be a short phrase (max 3-4 words).
    For each step, recommend a free online course that covers the topic in the context of {domain_name}.
    Include the course title and platform or website where it's available.
    
    Format the output as follows:
    1. Step Name: Course Title (Platform/Website)
    2. Step Name: Course Title (Platform/Website)
    ...
    
    Ensure diversity in course recommendations, suggesting different platforms and courses for each step when possible.
    """

    # Invoke the agent to generate the roadmap and course recommendations
    response = agent_executor.invoke({"input": combined_prompt})

    # Process the response to extract steps and course recommendations
    lines = response['output'].split('\n')
    roadmap_steps = []
    course_recommendations = []

    for line in lines:
        if line.strip():
            parts = line.split(':', 1)
            if len(parts) == 2:
                step = parts[0].strip()
                course = parts[1].strip()
                roadmap_steps.append(step)
                course_recommendations.append(course)

    # Create a graph using NetworkX
    G = nx.DiGraph()

    # Adding nodes and edges based on the roadmap steps
    for i, step in enumerate(roadmap_steps):
        G.add_node(i, label=f"{step}")
        if i > 0:
            G.add_edge(i - 1, i)

    # Set up the plot
    plt.figure(figsize=(16, 10))

    # Use a more compact layout
    pos = nx.spring_layout(G, k=0.8, iterations=50)

    # Draw the graph
    nx.draw(G, pos, node_shape='s', node_size=4000, node_color='lightblue', 
            edgecolors='black', linewidths=2, with_labels=False)

    # Add labels to nodes
    labels = nx.get_node_attributes(G, 'label')
    nx.draw_networkx_labels(G, pos, labels, font_size=10, font_weight='bold')

    # Add edge arrows
    nx.draw_networkx_edges(G, pos, edge_color='gray', arrows=True, arrowsize=20)

    # Set title and remove axes
    plt.title(f"{domain_name} Learning Roadmap", fontsize=20, fontweight='bold')
    plt.axis('off')

    # Adjust layout and save the figure
    output_path = f"{domain_name.lower().replace(' ', '_')}_roadmap.png"
    plt.tight_layout()
    plt.savefig(output_path, format='png', dpi=300, bbox_inches='tight')

    # Return the course recommendations and the image path
    return {
        "roadmap_image": output_path,
        "course_recommendations": {step: course for step, course in zip(roadmap_steps, course_recommendations)}
    }

@app.get("/roadmap-image/{domain_name}")
def get_roadmap_image(domain_name: str):
    output_path = f"{domain_name.lower().replace(' ', '_')}_roadmap.png"
    if os.path.exists(output_path):
        return FileResponse(output_path)
    raise HTTPException(status_code=404, detail="Roadmap image not found.")



if __name__ == "__main__":
    uvicorn.run(app, host="127.0.0.1", port=9000)
